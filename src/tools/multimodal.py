"""Turn uploaded files into Gemini multimodal Parts for the chat agent.

Gemini 2.5 natively analyses images, audio, video and PDF/text — so the agent reads the
file as real context (not just a filename). Small files go INLINE (``Part.from_bytes``);
files above the inline request cap go through the Gemini FILES API (``Part.from_uri``)
using the same ``GEMINI_API_KEY`` the agent's model uses, so the model can access them.
Plain text is decoded inline; unsupported binaries (e.g. .docx) degrade to a labelled
note instead of breaking the turn.

Pure helpers (``resolve_mime`` / ``classify``) are unit-tested; ``build_file_part`` does
the (async) upload and is exercised in the live smoke.
"""
from __future__ import annotations

import asyncio
import io
import os
from typing import Optional

from google.genai import types

INLINE_MAX = 18 * 1024 * 1024     # ~18 MB — under Gemini's ~20 MB inline request cap
UPLOAD_MAX = 100 * 1024 * 1024    # 100 MB — via the Files API (video / long audio)

# Extension → mime fallback when the browser omits content_type (or sends octet-stream).
_EXT_MIME = {
    ".pdf": "application/pdf", ".txt": "text/plain", ".csv": "text/csv", ".md": "text/markdown",
    ".mp3": "audio/mpeg", ".wav": "audio/wav", ".ogg": "audio/ogg", ".oga": "audio/ogg",
    ".m4a": "audio/mp4", ".aac": "audio/aac", ".flac": "audio/flac", ".weba": "audio/webm",
    ".mp4": "video/mp4", ".mov": "video/quicktime", ".webm": "video/webm",
    ".avi": "video/x-msvideo", ".mpeg": "video/mpeg", ".mpg": "video/mpeg", ".3gp": "video/3gpp",
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
    ".gif": "image/gif", ".webp": "image/webp", ".heic": "image/heic",
}
_NATIVE_PREFIXES = ("image/", "audio/", "video/")
_NATIVE_DOCS = {"application/pdf"}
_TEXT_PREFIX = "text/"


def resolve_mime(filename: str, content_type: Optional[str]) -> str:
    """Best mime for a file: the declared content_type, else the extension fallback."""
    ct = (content_type or "").split(";")[0].strip().lower()
    if ct and ct != "application/octet-stream":
        return ct
    ext = os.path.splitext(filename or "")[1].lower()
    return _EXT_MIME.get(ext, ct or "application/octet-stream")


def classify(mime: str) -> str:
    """'media' (image/audio/video/pdf → native Part), 'text' (decode inline), or
    'unsupported' (note only)."""
    if mime.startswith(_NATIVE_PREFIXES) or mime in _NATIVE_DOCS:
        return "media"
    if mime.startswith(_TEXT_PREFIX):
        return "text"
    return "unsupported"


def _gemini_client():
    from google import genai

    key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GEMINI_API_KEY_1")
    return genai.Client(api_key=key) if key else None


async def _upload_to_files_api(content: bytes, mime: str, filename: str) -> Optional[types.Part]:
    """Upload a large file to the Gemini Files API and wait until it's ACTIVE. The genai
    client is sync → run it off the event loop; poll with async sleeps."""
    client = _gemini_client()
    if client is None:
        return None
    up = await asyncio.to_thread(
        client.files.upload,
        file=io.BytesIO(content),
        config={"mime_type": mime, "display_name": filename or "file"},
    )
    for _ in range(45):  # video/large files need a few seconds to process
        f = await asyncio.to_thread(client.files.get, name=up.name)
        state = str(getattr(f, "state", "") or "")
        if state.endswith("ACTIVE"):
            return types.Part.from_uri(file_uri=up.uri, mime_type=mime)
        if state.endswith("FAILED"):
            return None
        await asyncio.sleep(1)
    return None  # timed out waiting for processing


async def analyze_attachments(parts: list, user_message: str = "",
                              model: str = "gemini-2.5-flash") -> str:
    """Analyse attached media DIRECTLY with Gemini (the path verified to actually read the
    file), with a business-aware system prompt. We route multimodal turns here instead of
    the multi-agent Runner, which drops non-text parts somewhere in its session/router
    pipeline. ``parts`` already contains the user's text Part + the media Parts."""
    client = _gemini_client()
    if client is None:
        return "No hay un cliente Gemini configurado para analizar el archivo."
    system = types.Part(text=(
        "Sos ARIA, el asistente de negocios del usuario (PyME). Analizá el/los archivo(s) "
        "adjunto(s) — audio, video, documento, imagen — y respondé de forma clara, útil y "
        "accionable para su negocio. Si el usuario hizo una pregunta, respondéla usando el "
        "contenido del archivo. Si es un audio/video, resumí lo relevante."
    ))
    contents = [system, *parts]
    try:
        resp = await asyncio.to_thread(client.models.generate_content, model=model, contents=contents)
        return (getattr(resp, "text", "") or "").strip() or "No pude extraer una respuesta del archivo."
    except Exception as e:  # noqa: BLE001
        return f"No pude analizar el archivo: {e}"


def _excel_to_text(content: bytes, max_rows: int = 300) -> str:
    """Flatten an .xlsx/.xls workbook to text (all sheets) so the model can read it."""
    import pandas as pd

    sheets = pd.read_excel(io.BytesIO(content), sheet_name=None, dtype=str)
    out = []
    for name, df in sheets.items():
        df = df.fillna("")
        lines = [" | ".join(str(c) for c in df.columns)]
        for _, row in df.head(max_rows).iterrows():
            lines.append(" | ".join(str(v) for v in row.values))
        out.append(f"--- Hoja '{name}' ({len(df)} filas) ---\n" + "\n".join(lines))
    return "\n\n".join(out)[:200_000]


def _pptx_to_text(content: bytes) -> str:
    """Extract slide text from a .pptx presentation."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        texts.append(t)
        if texts:
            out.append(f"--- Diapositiva {i} ---\n" + "\n".join(texts))
    return "\n\n".join(out)[:200_000]


def _docx_to_text(content: bytes) -> str:
    """Extract paragraph + table text from a .docx document."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)[:200_000]


def _extract_office(content: bytes, mime: str, filename: str) -> Optional[str]:
    """Excel / PowerPoint / Word → text (Gemini doesn't read those binaries natively).
    None if not an office file (falls through to the unsupported note)."""
    name = (filename or "").lower()
    is_xls = "spreadsheet" in mime or "ms-excel" in mime or name.endswith((".xlsx", ".xls"))
    is_ppt = "presentation" in mime or "powerpoint" in mime or name.endswith((".pptx", ".ppt"))
    is_doc = "wordprocessingml" in mime or "msword" in mime or name.endswith((".docx", ".doc"))
    try:
        if is_xls:
            return _excel_to_text(content)
        if is_ppt:
            return _pptx_to_text(content)
        if is_doc:
            return _docx_to_text(content)
    except Exception:  # noqa: BLE001 — corrupt/odd office file → fall back to a note
        return None
    return None


async def build_file_part(content: bytes, filename: str, content_type: Optional[str]
                          ) -> tuple[Optional[types.Part], str]:
    """One uploaded file → (Gemini Part | None, human-readable note).

    Native media is analysed by the model (inline, or Files API when large); text is
    decoded inline; unsupported binaries return ``None`` + a note."""
    mime = resolve_mime(filename, content_type)
    size = len(content)
    name = filename or "archivo"
    kind = classify(mime)

    if kind == "text":
        try:
            text = content.decode("utf-8", errors="replace")[:200_000]
        except Exception:  # noqa: BLE001
            text = ""
        return types.Part(text=f"[Contenido de {name} ({mime})]:\n{text}"), f"[{name} · texto]"

    if kind == "media":
        if size <= INLINE_MAX:
            return types.Part.from_bytes(data=content, mime_type=mime), f"[{name} · {mime}]"
        if size <= UPLOAD_MAX:
            part = await _upload_to_files_api(content, mime, name)
            if part is not None:
                return part, f"[{name} · {mime} · {size // 1024 // 1024}MB vía Files API]"
            return None, f"[{name}: el archivo grande no pudo procesarse]"
        return None, f"[{name}: excede el máximo de {UPLOAD_MAX // 1024 // 1024}MB]"

    office = await asyncio.to_thread(_extract_office, content, mime, name)
    if office:
        return types.Part(text=f"[Contenido de {name}]:\n{office}"), f"[{name} · Office → texto extraído]"
    return None, f"[Archivo adjunto no analizable directamente: {name} ({mime}, {size} bytes)]"
