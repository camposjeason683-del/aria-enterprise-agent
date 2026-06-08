"""Unit tests for the multimodal helpers (mime resolution + classification + Office text
extraction). The actual Gemini upload is exercised in the live smoke."""
import io

from src.tools.multimodal import (
    _docx_to_text, _excel_to_text, _extract_office, _pptx_to_text, classify, resolve_mime,
)


def test_resolve_mime_prefers_declared_content_type():
    assert resolve_mime("x.bin", "audio/mpeg") == "audio/mpeg"


def test_resolve_mime_extension_fallback_when_missing_or_octet():
    assert resolve_mime("nota.mp3", "application/octet-stream") == "audio/mpeg"
    assert resolve_mime("clip.mp4", None) == "video/mp4"
    assert resolve_mime("informe.pdf", "") == "application/pdf"
    assert resolve_mime("grab.m4a", None) == "audio/mp4"
    assert resolve_mime("x.desconocido", None) == "application/octet-stream"


def test_resolve_mime_strips_charset():
    assert resolve_mime("a.txt", "text/plain; charset=utf-8") == "text/plain"


def test_classify_media_text_unsupported():
    for m in ("image/png", "audio/mpeg", "video/mp4", "video/quicktime", "application/pdf"):
        assert classify(m) == "media"
    assert classify("text/plain") == "text"
    assert classify("text/csv") == "text"
    # office binary → not native to Gemini (extracted to text separately, see below)
    assert classify("application/vnd.openxmlformats-officedocument.wordprocessingml.document") == "unsupported"


def test_excel_to_text():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["producto", "ventas"])
    ws.append(["Tomate", "5000"])
    buf = io.BytesIO()
    wb.save(buf)
    txt = _excel_to_text(buf.getvalue())
    assert "producto" in txt and "Tomate" in txt and "5000" in txt


def test_pptx_to_text():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "Ventas Q2: 12000"
    buf = io.BytesIO()
    prs.save(buf)
    txt = _pptx_to_text(buf.getvalue())
    assert "12000" in txt


def test_docx_to_text():
    from docx import Document

    doc = Document()
    doc.add_paragraph("Contrato con proveedor X")
    doc.add_paragraph("Monto total: 45000 pesos")
    buf = io.BytesIO()
    doc.save(buf)
    txt = _docx_to_text(buf.getvalue())
    assert "45000" in txt and "proveedor" in txt


def test_extract_office_routes_by_mime_and_ext():
    import openpyxl

    wb = openpyxl.Workbook()
    wb.active.append(["x", "y"])
    buf = io.BytesIO()
    wb.save(buf)
    # routes by extension even when mime is octet-stream (browser omitted it)
    assert _extract_office(buf.getvalue(), "application/octet-stream", "datos.xlsx") is not None
    # a non-office file → None (falls through to the unsupported note)
    assert _extract_office(b"hola", "text/plain", "n.txt") is None
